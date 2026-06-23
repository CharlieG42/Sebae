# This Python file uses the following encoding: utf-8
"""
Génération de rapports PowerPoint.
Fusionne ECReport (MyApp) et PrintReport (WUEnergyTool).

Principe :
  - replaceTxtInPptx : remplace les balises $tag$ par leur valeur
  - replaceImgInPptx : remplace une balise $tag$ par une image
  - print_pptx       : génère le rapport complet EC
"""

from __future__ import annotations

import os
import sys
import locale
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm as pptxCm

try:
    locale.setlocale(locale.LC_TIME, "fr_FR.UTF-8")
except locale.Error:
    pass   # Windows peut ne pas avoir ce locale — on continue sans


class PrintReport:

    def __init__(self):
        if getattr(sys, "frozen", False):
            self.base_dir = Path(sys.executable).resolve().parent
        else:
            self.base_dir = Path(__file__).resolve().parent.parent

    # -------------------------------------------------------------------------

    def replace_txt_in_pptx(self, prs: Presentation, replacements: list[list[str]]) -> bool:
        """
        Remplace chaque balise [tag, value] dans tous les runs du fichier pptx.
        Retourne True si au moins un remplacement a eu lieu.
        """
        for tag, value in replacements:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if tag in run.text:
                                    run.text = run.text.replace(tag, str(value))
        return True

    def replace_img_in_pptx(
        self,
        prs:       Presentation,
        tag:       str,
        img_path:  str,
        width_cm:  float = 22.0,
        height_cm: float = 15.0,
    ) -> bool:
        """
        Remplace la forme contenant 'tag' par une image positionnée au même endroit.
        """
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and tag in shape.text:
                    left  = shape.left
                    top   = shape.top
                    shape._element.getparent().remove(shape._element)
                    slide.shapes.add_picture(
                        img_path, left, top,
                        width=pptxCm(width_cm),
                        height=pptxCm(height_cm),
                    )
                    break   # une seule occurrence par slide
        return True

    def print_pptx(self, data: dict, template_path: str | None = None) -> Path:
        """
        Génère le rapport Energy Check à partir d'un template pptx.

        Paramètres
        ----------
        data          : dictionnaire de valeurs (clés = noms de balises sans $)
        template_path : chemin vers le template (défaut : assets/template_EC.pptx)

        Retourne le chemin du fichier généré.
        """
        if template_path is None:
            template_path = str(self.base_dir / "assets" / "template_EC.pptx")

        prs   = Presentation(template_path)
        today = datetime.now().strftime("%d %B %Y")

        replacements = [
            ["$title$",       data.get("title",       "")],
            ["$customer$",    data.get("customer",     "")],
            ["$eco_euros$",   data.get("eco_euros",    "")],
            ["$roi$",         data.get("roi",          "")],
            ["$eco_nrj$",     data.get("eco_nrj",      "")],
            ["$eco_co2$",     data.get("eco_co2",      "")],
            ["$invest$",      data.get("invest",       "")],
            ["$eco_10y$",     data.get("eco_10y",      "")],
            ["$eco_15y$",     data.get("eco_15y",      "")],
            ["$trees_qty$",   data.get("trees_qty",    "")],
            ["$cars_qty$",    data.get("cars_qty",     "")],
            ["$nrj_cost$",    data.get("nrj_cost",     "")],
            ["$total_invest$",data.get("total_invest", "")],
            ["$date$",        today],
            ["$eff_loss$",    data.get("effLoss",      "")],
            ["$pump1$",       data.get("pump1",        "")],
            ["$pump2$",       data.get("pump2",        "")],
            ["$q_h1$",        data.get("q_h1",         "")],
            ["$q_h2$",        data.get("q_h2",         "")],
            ["$time1$",       data.get("time1",        "")],
            ["$time2$",       data.get("time2",        "")],
            ["$power1$",      data.get("power1",       "")],
            ["$power2$",      data.get("power2",       "")],
            ["$eff1$",        data.get("eff1",         "")],
            ["$eff2$",        data.get("eff2",         "")],
        ]

        self.replace_txt_in_pptx(prs, replacements)
        self.replace_img_in_pptx(prs, "$picture1$", data.get("img_cost",  "graphLineCost.png"))
        self.replace_img_in_pptx(prs, "$picture2$", data.get("img_power", "graphLinePower.png"))

        output_dir = self.base_dir / "output"
        output_dir.mkdir(exist_ok=True)
        output_path = output_dir / f'{data.get("title", "rapport")}.pptx'

        prs.save(output_path)

        # Ouverture automatique sur Windows
        if sys.platform == "win32":
            os.startfile(output_path)

        return output_path

    # --- Alias pour compatibilité avec l'ancien code ---
    def replaceTxtInPptx(self, prs, lst):   return self.replace_txt_in_pptx(prs, lst)
    def replaceImgInPptx(self, prs, t, i):  return self.replace_img_in_pptx(prs, t, i)
    def printPptx(self, data):              return self.print_pptx(data)
